"""统一文档文本抽取：PDF（含扫描件 OCR 兜底）/ DOCX / PPTX / XLSX / TXT / MD / HTML / CSV。

- PDF 用 PyMuPDF 提文字；若某页几乎无文字（扫描件），渲染成图走 RapidOCR（本地、中英文）。
- OCR 引擎懒加载（首次使用才初始化，避免拖慢启动）。
"""
import os
import re

_ocr = None


def _get_ocr():
    global _ocr
    if _ocr is None:
        from rapidocr_onnxruntime import RapidOCR
        _ocr = RapidOCR()
    return _ocr


def _ocr_image_bytes(png_bytes: bytes) -> str:
    try:
        import numpy as np
        from PIL import Image
        import io
        img = np.array(Image.open(io.BytesIO(png_bytes)).convert("RGB"))
        result, _ = _get_ocr()(img)
        if not result:
            return ""
        return "\n".join(line[1] for line in result)
    except Exception as e:
        print(f"[extract] OCR failed: {e}")
        return ""


def _extract_pdf(path: str) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(path)
    pages_text = []
    scanned_pages = []
    for i, page in enumerate(doc):
        t = page.get_text("text") or ""
        pages_text.append(t)
        if len(t.strip()) < 40:  # 该页几乎无文字 → 可能是扫描图
            scanned_pages.append(i)

    total_chars = sum(len(t) for t in pages_text)
    # 整篇文字极少 / 大量空页 → 走 OCR（限制页数，避免超大扫描件过慢）
    needs_ocr = total_chars < 200 or len(scanned_pages) > len(pages_text) * 0.5
    if needs_ocr and scanned_pages:
        print(f"[extract] PDF looks scanned ({len(scanned_pages)}/{len(pages_text)} pages), running OCR…")
        ocr_pages = scanned_pages[:40]  # 最多 OCR 40 页
        for i in ocr_pages:
            pix = doc[i].get_pixmap(dpi=200)
            pages_text[i] = (pages_text[i] + "\n" + _ocr_image_bytes(pix.tobytes("png"))).strip()
    doc.close()
    return "\n".join(pages_text)


def _extract_docx(path: str) -> str:
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text]
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def _extract_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"--- Sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _extract_text_file(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        txt = f.read()
    if path.lower().endswith((".html", ".htm")):
        txt = re.sub(r"<[^>]+>", " ", txt)  # 粗略去标签
    return txt


_DISPATCH = {
    ".pdf": _extract_pdf, ".docx": _extract_docx, ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx, ".xls": _extract_xlsx,
    ".txt": _extract_text_file, ".md": _extract_text_file, ".csv": _extract_text_file,
    ".html": _extract_text_file, ".htm": _extract_text_file,
}

SUPPORTED_EXTENSIONS = sorted(_DISPATCH.keys())


def extract_text(path: str) -> str:
    """按扩展名抽取文档全文。未知类型按 UTF-8 文本兜底。"""
    ext = os.path.splitext(path)[1].lower()
    fn = _DISPATCH.get(ext, _extract_text_file)
    return fn(path)
